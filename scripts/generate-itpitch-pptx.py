#!/usr/bin/env python3
"""Generate a slide-deck version of the Dr. Asa IT Pitch."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = os.path.join(os.path.dirname(__file__), '..', 'docs', 'deliverables',
                   'UNFPA-IT-Pitch-DrAsa.pptx')

DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
LIGHT_BLUE = RGBColor(0x00, 0x9E, 0xDB)
PALE_BLUE = RGBColor(0xF0, 0xF6, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = LIGHT_BLUE
        shp.line.width = Pt(1.5)
    return shp


def tb(slide, left, top, width, height, text, size=18, bold=False,
       color=DARK_GRAY, align=PP_ALIGN.LEFT, font='Calibri'):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def bullets(slide, left, top, width, height, items, size=16,
            color=DARK_GRAY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = '•  '
            r1.font.size = Pt(size)
            r1.font.color.rgb = LIGHT_BLUE
            r1.font.bold = True
            r1.font.name = 'Calibri'
            r2 = p.add_run()
            r2.text = lead
            r2.font.size = Pt(size)
            r2.font.bold = True
            r2.font.color.rgb = DARK_BLUE
            r2.font.name = 'Calibri'
            r3 = p.add_run()
            r3.text = rest
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
            r3.font.name = 'Calibri'
        else:
            r1 = p.add_run()
            r1.text = '•  '
            r1.font.size = Pt(size)
            r1.font.color.rgb = LIGHT_BLUE
            r1.font.bold = True
            r1.font.name = 'Calibri'
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = 'Calibri'
    return tf


def header(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.05),
         DARK_BLUE)
    rect(slide, Inches(0), Inches(1.05), prs.slide_width, Inches(0.08),
         LIGHT_BLUE)
    tb(slide, Inches(0.6), Inches(0.22), Inches(12), Inches(0.7),
       title, size=28, bold=True, color=WHITE)
    if subtitle:
        tb(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.35),
           subtitle, size=13, color=RGBColor(0xBB, 0xDD, 0xFF))


def footer(slide, page, total):
    tb(slide, Inches(0.6), Inches(7.05), Inches(10), Inches(0.35),
       'UNFPA Partnership Catalyst  |  IT Team Pitch  |  Dr. Asa Torkelsson',
       size=10, color=MED_GRAY)
    tb(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
       f'{page} / {total}', size=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)


slides = []

# ── 1: Cover ──
def slide1():
    s = prs.slides.add_slide(BLANK)
    bg(s, DARK_BLUE)
    rect(s, Inches(0), Inches(5.3), prs.slide_width, Inches(0.1), LIGHT_BLUE)
    tb(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.5),
       'UNFPA PARTNERSHIP CATALYST',
       size=18, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.3),
       'Introducing an AI-Powered\nKnowledge Tool',
       size=44, bold=True, color=WHITE)
    tb(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
       'A briefing for the UNFPA IT Team',
       size=20, color=RGBColor(0xBB, 0xDD, 0xFF))
    tb(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.5),
       'Dr. Asa Torkelsson',
       size=18, bold=True, color=WHITE)
    tb(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.5),
       'Chief, UNFPA Seoul Representation Office',
       size=14, color=RGBColor(0xBB, 0xBB, 0xBB))
    tb(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
       'April 2026',
       size=12, color=RGBColor(0x99, 0x99, 0x99))
slides.append(slide1)

# ── 2: The Problem ──
def slide2():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'The Problem: Two Broken Sides',
           'A structural mismatch between capacity and capital')
    # UNFPA column
    rect(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(4.8), PALE_BLUE,
         line=True)
    tb(s, Inches(0.85), Inches(1.8), Inches(5.5), Inches(0.5),
       'UNFPA SIDE', size=14, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(0.85), Inches(2.2), Inches(5.5), Inches(0.6),
       'Has problems + solutions',
       size=20, bold=True, color=DARK_BLUE)
    bullets(s, Inches(0.85), Inches(2.9), Inches(5.5), Inches(3.2), [
        'Proven programmes across 150+ countries',
        'Maternal health, SRHR, GBV, midwifery delivery',
        'Deep country presence and government trust',
        'Cannot access non-traditional funding at scale',
    ], size=14)
    # Philanthropy column
    rect(s, Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.8), PALE_BLUE,
         line=True)
    tb(s, Inches(7.05), Inches(1.8), Inches(5.5), Inches(0.5),
       'PHILANTHROPY SIDE', size=14, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(7.05), Inches(2.2), Inches(5.5), Inches(0.6),
       'Has capital ready to deploy',
       size=20, bold=True, color=DARK_BLUE)
    bullets(s, Inches(7.05), Inches(2.9), Inches(5.5), Inches(3.2), [
        'Hundreds of family offices and foundations',
        'Active interest in health, climate, gender',
        'Strong Singapore and Asia-Pacific presence',
        'Does not know where or how to invest',
    ], size=14)
    tb(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
       'Two worlds that should be finding each other are passing each other by.',
       size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
slides.append(slide2)

# ── 3: The Reframe ──
def slide3():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'The Strategic Reframe',
           'Not research. Not a report. A market-maker.')
    tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
       'This tool is:',
       size=18, bold=True, color=DARK_BLUE)
    tb(s, Inches(0.9), Inches(2.15), Inches(12), Inches(0.5),
       '✗   NOT just a PPP design exercise',
       size=18, color=MED_GRAY)
    tb(s, Inches(0.9), Inches(2.65), Inches(12), Inches(0.5),
       '✗   NOT just a static research report',
       size=18, color=MED_GRAY)
    rect(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(2.4), PALE_BLUE,
         line=True)
    tb(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.5),
       '✓   A MARKET-CREATION AND MATCHMAKING SYSTEM',
       size=16, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.6),
       'A live intelligence platform that bridges the two broken sides — '
       'giving UNFPA staff the intelligence and language they need to '
       'initiate, prepare for, and close funding conversations with '
       'non-traditional partners.',
       size=16, color=DARK_GRAY)
    tb(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.5),
       'The real question: do we participate in the philanthropic market '
       'on our own terms, or wait to be discovered?',
       size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
slides.append(slide3)

# ── 4: LKYSPP Deliverables ──
def slide4():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'What the LKYSPP Team Delivered',
           'Two concrete, production-ready deliverables')
    # Deliverable 1
    rect(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5.1), PALE_BLUE,
         line=True)
    tb(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(0.4),
       'DELIVERABLE 1', size=12, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(0.85), Inches(2.1), Inches(5.5), Inches(0.8),
       'Singapore Philanthropy\nLandscape Report',
       size=22, bold=True, color=DARK_BLUE)
    tb(s, Inches(0.85), Inches(3.4), Inches(5.5), Inches(0.4),
       '60 pages of original research',
       size=13, bold=True, color=LIGHT_BLUE)
    bullets(s, Inches(0.85), Inches(3.85), Inches(5.5), Inches(2.6), [
        'Map of family offices and foundations',
        'Funder priorities and decision criteria',
        'Alignment points with UNFPA mandate',
        'Embedded and queryable inside the tool',
    ], size=13)
    # Deliverable 2
    rect(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.1), PALE_BLUE,
         line=True)
    tb(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4),
       'DELIVERABLE 2', size=12, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.8),
       'AI-Powered Knowledge\nTool (Prototype)',
       size=22, bold=True, color=DARK_BLUE)
    tb(s, Inches(7.05), Inches(3.4), Inches(5.5), Inches(0.4),
       'Live web application',
       size=13, bold=True, color=LIGHT_BLUE)
    bullets(s, Inches(7.05), Inches(3.85), Inches(5.5), Inches(2.6), [
        'Curated UNFPA knowledge base (32 docs)',
        'RAG chat for partnership preparation',
        'Export to DOCX, PDF, PPTX',
        'Already running: unfpa-lkyspp-otg.vercel.app',
    ], size=13)
slides.append(slide4)

# ── 5: What the Tool Can Do ──
def slide5():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'What the AI Tool Can Do',
           'Six core partnership-preparation capabilities')
    bullets(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.5), [
        ('Generate funder pitches  ',
         '— framing UNFPA programmes for specific family offices and foundations'),
        ('Draft briefing notes  ',
         '— structured two-pagers ready for external partners'),
        ('Prepare meeting talking points  ',
         '— discussion frameworks for funding conversations'),
        ('Match projects to funders  ',
         '— identify alignment between programmes and donor priorities'),
        ('Frame work for climate funding  ',
         '— position SRHR within climate resilience narratives'),
        ('Compare financing models  ',
         '— DIBs, blended finance, South-South cooperation, explained'),
    ], size=17)
slides.append(slide5)

# ── 6: Under the Hood ──
def slide6():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'Under the Hood',
           'Production-grade, standard, no vendor lock-in')
    tb(s, Inches(0.6), Inches(1.55), Inches(6), Inches(0.4),
       'TECHNOLOGY STACK', size=13, bold=True, color=LIGHT_BLUE)
    bullets(s, Inches(0.6), Inches(2.0), Inches(6.2), Inches(4.8), [
        ('Frontend  ', '— Next.js 16, React 19, TypeScript, Tailwind'),
        ('Backend  ', '— Node.js API routes'),
        ('Database  ', '— PostgreSQL + pgvector extension'),
        ('ORM  ', '— Prisma 7'),
        ('Language model  ', '— Anthropic Claude Sonnet 4'),
        ('Embeddings  ', '— OpenAI text-embedding-3-small'),
        ('Hosting  ', '— Vercel (currently), portable to any cloud'),
    ], size=13)
    tb(s, Inches(7.1), Inches(1.55), Inches(6), Inches(0.4),
       'HOW IT WORKS', size=13, bold=True, color=LIGHT_BLUE)
    bullets(s, Inches(7.1), Inches(2.0), Inches(6), Inches(4.8), [
        'Markdown knowledge docs ingested via CLI script',
        'SHA-256 hash prevents re-ingesting unchanged docs',
        'Chunked at ~1000 words on heading boundaries',
        'Each chunk embedded into a 1,536-dim vector',
        'Stored in PostgreSQL with pgvector',
        'Queries embedded and matched via cosine similarity',
        'Claude synthesises retrieved chunks into a response',
    ], size=13)
slides.append(slide6)

# ── 7: Handover Terms ──
def slide7():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'Handover Terms',
           'A gift from LKYSPP, no strings attached')
    rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.0), PALE_BLUE)
    tb(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.7),
       'Cost of acquisition: ZERO',
       size=24, bold=True, color=DARK_BLUE)
    tb(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.4),
       'No licensing fee. No transfer fee. No ongoing obligation.',
       size=13, color=DARK_GRAY)
    bullets(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(4), [
        ('Codebase  ',
         '— full source transferred to private UNFPA GitHub organisation'),
        ('Hosting  ',
         '— Vercel free tier, or self-hosted on any cloud / on-premise'),
        ('Runtime APIs  ',
         '— Anthropic + OpenAI, ~USD 35–170/month for 20–30 users'),
        ('Knowledge stewardship  ',
         '— programme staff edit Markdown, run one command to reingest'),
        ('Technical stewardship  ',
         '— UNFPA IT team owns deployment, API keys, code changes'),
        ('Long-term maintenance  ',
         '— LKYSPP concludes at handover; standard stack, any contractor can maintain'),
    ], size=14)
slides.append(slide7)

# ── 8: IT Approvals Overview ──
def slide8():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'Questions for the IT Team',
           '30 questions across 6 governance areas')
    areas = [
        ('A. Data Governance & Privacy',
         'Third-party AI APIs, data residency, PII, DPO sign-off'),
        ('B. Cybersecurity & Vendor Risk',
         'Vendor assessment, Anthropic/OpenAI approval, pen test, cloud standards'),
        ('C. Procurement & Legal',
         'Spend thresholds, IP transfer instrument, Legal sign-off'),
        ('D. Infrastructure & Hosting',
         'Vercel acceptability, self-hosting, PostgreSQL + pgvector provisioning'),
        ('E. Access Control & Identity',
         'SSO integration, role-based access, audit logging'),
        ('F. AI & Algorithmic Governance',
         'Internal AI policy, human-in-the-loop review, AI inventory'),
    ]
    for i, (label, desc) in enumerate(areas):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.15)
        top = Inches(1.65 + row * 1.65)
        rect(s, left, top, Inches(5.95), Inches(1.5), PALE_BLUE, line=True)
        tb(s, left + Inches(0.2), top + Inches(0.15), Inches(5.6),
           Inches(0.4), label, size=14, bold=True, color=DARK_BLUE)
        tb(s, left + Inches(0.2), top + Inches(0.55), Inches(5.6),
           Inches(0.9), desc, size=12, color=DARK_GRAY)
slides.append(slide8)

# ── 9: Sample Approval Questions ──
def slide9():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'A Sample of the 30 Questions',
           'The full list lives in the accompanying document')
    bullets(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.5), [
        'Does UNFPA HQ policy permit third-party AI APIs to process internal documents?',
        'Are Anthropic and OpenAI on UNFPA\'s approved vendor list?',
        'Can UNFPA formally accept Vercel as a host for a staff-facing tool?',
        'Does the codebase transfer from LKYSPP require a legal instrument (MOU, IP)?',
        'Should the tool integrate with UNFPA SSO / Azure AD before any rollout?',
        'Does UNFPA have an internal AI use policy this tool must comply with?',
        'Is audit logging of chat queries required for compliance?',
        'Who in UNFPA Legal needs to sign off on accepting the codebase?',
    ], size=15)
slides.append(slide9)

# ── 10: My Ask ──
def slide10():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, 'My Ask of the IT Team',
           'A 30-day assessment window')
    tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
       'I am not asking you to build anything. The LKYSPP team has already built it.',
       size=17, bold=True, color=DARK_BLUE)
    tb(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.4),
       'I am asking you to:',
       size=15, color=DARK_GRAY)
    bullets(s, Inches(0.6), Inches(2.65), Inches(12.1), Inches(4.2), [
        ('Review  ', 'the codebase against UNFPA security and hosting standards'),
        ('Work through  ', 'the 30-question approval checklist'),
        ('Advise  ', 'on hosting (Vercel vs. self-hosted) and authentication requirements'),
        ('Facilitate  ', 'API vendor approval for Anthropic and OpenAI'),
        ('Participate  ', 'in a technical handover session with the LKYSPP development team'),
        ('Deliver  ', 'a written recommendation at the end of 30 days: deploy, modify, or escalate'),
    ], size=16)
slides.append(slide10)

# ── 11: Closing ──
def slide11():
    s = prs.slides.add_slide(BLANK)
    bg(s, DARK_BLUE)
    rect(s, Inches(0), Inches(3.7), prs.slide_width, Inches(0.08), LIGHT_BLUE)
    tb(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.5),
       'THE BOTTOM LINE', size=16, bold=True, color=LIGHT_BLUE)
    tb(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.7),
       'The market failure between UNFPA\'s programmatic capacity\n'
       'and philanthropic capital will not fix itself.',
       size=24, bold=True, color=WHITE)
    tb(s, Inches(0.9), Inches(4.0), Inches(11.5), Inches(1.2),
       'We need better tools to participate in that market —\n'
       'actively, intelligently, and on our own terms.',
       size=20, color=RGBColor(0xBB, 0xDD, 0xFF))
    tb(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.4),
       'Dr. Asa Torkelsson',
       size=16, bold=True, color=WHITE)
    tb(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.4),
       'Chief, UNFPA Seoul Representation Office',
       size=13, color=RGBColor(0xBB, 0xBB, 0xBB))
    tb(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
       'April 2026',
       size=11, color=RGBColor(0x99, 0x99, 0x99))
slides.append(slide11)


for fn in slides:
    fn()

total = len(prs.slides)
# Add footers (skip cover + closing)
for i, slide in enumerate(prs.slides):
    if i == 0 or i == total - 1:
        continue
    footer(slide, i + 1, total)

prs.save(OUT)
size = os.path.getsize(OUT)
print(f'PPTX written: {OUT} ({size:,} bytes, {total} slides)')
